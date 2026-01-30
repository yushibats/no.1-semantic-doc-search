import base64
import logging
import os
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

import PyPDF2
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """文書処理サービス - PDF/PPT/PPTX/PNG/JPG/JPEGの解析"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'ppt': self._process_pptx,
            'pptx': self._process_pptx,
            'docx': self._process_docx,
            'txt': self._process_txt,
            'md': self._process_txt,  # Markdownもテキストとして処理
            'png': self._process_image,
            'jpg': self._process_image,
            'jpeg': self._process_image,
        }
    
    def process_document(self, file_path: str, filename: str) -> Dict[str, Any]:
        """
        文書を処理してテキストを抽出
        
        Args:
            file_path: ファイルパス
            filename: ファイル名
            
        Returns:
            処理結果(テキストチャンクとメタデータ)
        """
        try:
            # ファイル拡張子を取得
            file_ext = Path(filename).suffix.lower().lstrip('.')
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"サポートされていないファイル形式: {file_ext}")
            
            # ファイル処理
            processor = self.supported_formats[file_ext]
            result = processor(file_path, filename)
            
            return {
                "success": True,
                "filename": filename,
                "file_type": file_ext,
                "chunks": result.get("chunks", []),
                "metadata": result.get("metadata", {}),
                "page_count": result.get("page_count", 0)
            }
            
        except Exception as e:
            logger.error(f"文書処理エラー: {filename} - {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def _process_pdf(self, file_path: str, filename: str) -> Dict[str, Any]:
        """PDF処理"""
        chunks = []
        
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            page_count = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    # ページごとにチャンク化
                    chunks.append({
                        "page_number": page_num,
                        "text": text.strip(),
                        "chunk_id": f"{uuid.uuid4()}"
                    })
        
        return {
            "chunks": chunks,
            "page_count": page_count,
            "metadata": {
                "format": "PDF",
                "total_pages": page_count
            }
        }
    
    def _process_pptx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """PowerPoint処理"""
        chunks = []
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            
            # スライド内の全テキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            
            slide_text = "\n".join(text_parts).strip()
            
            if slide_text:
                chunks.append({
                    "page_number": slide_num,
                    "text": slide_text,
                    "chunk_id": f"{uuid.uuid4()}"
                })
        
        return {
            "chunks": chunks,
            "page_count": slide_count,
            "metadata": {
                "format": "PowerPoint",
                "total_slides": slide_count
            }
        }
    
    def _process_docx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Word文書処理"""
        chunks = []
        doc = DocxDocument(file_path)
        
        # 段落ごとにテキストを抽出
        paragraphs_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs_text.append(para.text.strip())
        
        # 適切なサイズにチャンク化（1000文字程度）
        current_chunk = []
        current_length = 0
        chunk_size = 1000
        
        for para in paragraphs_text:
            para_length = len(para)
            
            if current_length + para_length > chunk_size and current_chunk:
                # 現在のチャンクを保存
                chunks.append({
                    "page_number": len(chunks) + 1,
                    "text": "\n\n".join(current_chunk),
                    "chunk_id": f"{uuid.uuid4()}"
                })
                current_chunk = [para]
                current_length = para_length
            else:
                current_chunk.append(para)
                current_length += para_length
        
        # 最後のチャンクを保存
        if current_chunk:
            chunks.append({
                "page_number": len(chunks) + 1,
                "text": "\n\n".join(current_chunk),
                "chunk_id": f"{uuid.uuid4()}"
            })
        
        return {
            "chunks": chunks,
            "page_count": len(chunks),
            "metadata": {
                "format": "Word",
                "total_paragraphs": len(paragraphs_text)
            }
        }
    
    def _process_txt(self, file_path: str, filename: str) -> Dict[str, Any]:
        """テキストファイル処理"""
        chunks = []
        
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # 適切なサイズにチャンク化（1000文字程度）
        chunk_size = 1000
        text_chunks = []
        
        for i in range(0, len(content), chunk_size):
            chunk_text = content[i:i + chunk_size].strip()
            if chunk_text:
                text_chunks.append(chunk_text)
        
        for idx, chunk_text in enumerate(text_chunks, start=1):
            chunks.append({
                "page_number": idx,
                "text": chunk_text,
                "chunk_id": f"{uuid.uuid4()}"
            })
        
        return {
            "chunks": chunks,
            "page_count": len(chunks),
            "metadata": {
                "format": "Text",
                "total_chars": len(content)
            }
        }
    
    def _process_image(self, file_path: str, filename: str) -> Dict[str, Any]:
        """画像ファイル処理 - OCI Vision AIを使用してテキスト抽出"""
        try:
            from app.services.ai_copilot import get_copilot_service
            
            # 画像をbase64エンコード
            with open(file_path, 'rb') as f:
                image_data = f.read()
            
            # 画像の検証
            try:
                img = Image.open(file_path)
                width, height = img.size
                logger.info(f"画像サイズ: {width}x{height}")
            except Exception as e:
                logger.error(f"画像の検証エラー: {e}")
                raise ValueError(f"無効な画像ファイル: {e}")
            
            # base64エンコード
            base64_image = base64.b64encode(image_data).decode('utf-8')
            
            # MIMEタイプを判定
            file_ext = Path(filename).suffix.lower().lstrip('.')
            mime_type = f"image/{file_ext}" if file_ext in ['png', 'jpg', 'jpeg'] else 'image/jpeg'
            data_url = f"data:{mime_type};base64,{base64_image}"
            
            # Vision AIでテキスト抽出
            copilot = get_copilot_service()
            prompt = "この画像に含まれるすべてのテキストを抽出してください。テキストがない場合は「テキストなし」と応答してください。"
            
            # 同期的にテキスト抽出（ストリーミング不要）
            import asyncio
            
            async def extract_text():
                result = []
                async for chunk in copilot.chat_stream(
                    message=prompt,
                    images=[{"data_url": data_url}]
                ):
                    result.append(chunk)
                return ''.join(result)
            
            # イベントループを取得または作成
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
            
            extracted_text = loop.run_until_complete(extract_text())
            
            logger.info(f"画像からテキストを抽出: {len(extracted_text)}文字")
            
            # テキストがない場合
            if not extracted_text or extracted_text.strip().lower() in ['テキストなし', 'no text']:
                extracted_text = f"[画像: {filename}]\n画像からテキストを抽出できませんでした。"
            
            # チャンクとして返す
            chunks = [{
                "page_number": 1,
                "text": extracted_text.strip(),
                "chunk_id": f"{uuid.uuid4()}"
            }]
            
            return {
                "chunks": chunks,
                "page_count": 1,
                "metadata": {
                    "format": "Image",
                    "image_size": f"{width}x{height}",
                    "mime_type": mime_type
                }
            }
            
        except Exception as e:
            logger.error(f"画像処理エラー: {filename} - {e}")
            # エラーの場合でも基本情報を返す
            return {
                "chunks": [{
                    "page_number": 1,
                    "text": f"[画像: {filename}]\n画像処理エラー: {str(e)}",
                    "chunk_id": f"{uuid.uuid4()}"
                }],
                "page_count": 1,
                "metadata": {
                    "format": "Image",
                    "error": str(e)
                }
            }

# シングルトンインスタンス
document_processor = DocumentProcessor()
